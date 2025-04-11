from office365.sharepoint.client_context import ClientContext
from office365.runtime.auth.client_credential import ClientCredential
from services.llm_service import LLMService
from llama_index.core import Document
from typing import List
import tempfile
import pptx
import os

class SharePointLoader:
    def __init__(self, site_url: str, client_id: str, client_secret: str, folder_path: str, llm_service: LLMService):
        self.site_url = site_url
        self.client_id = client_id
        self.client_secret = client_secret
        self.llm_service = llm_service
        self.folder_path = folder_path
        self.ctx = self._authenticate()

    def _authenticate(self):
        credentials = ClientCredential(self.client_id, self.client_secret)
        return ClientContext(self.site_url).with_credentials(credentials)

    def load_documents(self) -> List[Document]:
        folder = self.ctx.web.get_folder_by_server_relative_url(self.folder_path)
        files = folder.expand(["Files"]).get().execute_query()
        documents = []

        for file in files.files:
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                file.download(tmp.name).execute_query()

                if file.name.endswith(('.mp4', '.mov', '.avi')):
                    documents.extend(self._process_video(tmp.name, file))
                elif file.name.endswith(('.pptx', '.ppt')):
                    documents.extend(self._process_slides(tmp.name, file))

        return documents

    def _process_video(self, file_path, file_item):
        transcript = self.llm_service.transcribe_audio(file_path)
        return [
            Document(
                text=segment['text'],
                metadata={
                    'type': 'video',
                    'title': file_item.name,
                    'timecodes': segment['start'],
                    'sharepoint_url': file_item.serverRelativeUrl
                }
            )
            for segment in transcript.segments
        ]

    def _process_slides(self, file_path, file_item):
        presentation = pptx.Presentation(file_path)
        return [
            Document(
                text=f"Slide {i + 1}: {shape.text}" if shape.has_text_frame else "",
                metadata={
                    'type': 'slide',
                    'title': file_item.name,
                    'slide_number': i + 1,
                    'sharepoint_url': file_item.serverRelativeUrl
                }
            )
            for i, slide in enumerate(presentation.slides)
            for shape in slide.shapes
            if shape.has_text_frame
        ]


#####################################################
# ORIGINAL CODE
# #########################################################
#
# import tempfile
# import tempfil
# import os
# import pptx
# from llama_index import Document
# from typing import List
# from office365.sharepoint.client_context import ClientContext
# from office365.runtime.auth.client_credential import ClientCredential
# from office365.runtime.auth.authentication_context import AuthenticationContext
# from config.config import Config
# from services.llm_service import LLMService

# class SharePointLoader:
#     def __init__(self, client, llm_service, temp_dir):
#         self.client = client
#         self.llm_service = llm_service
#         self.temp_dir = temp_dir
#         self.stats = {
#             'videos': 0,
#             'slides': 0,
#             'other': 0
#         }

#     def load_folder(self, folder_path) -> List[Document]:
#         folder = self.client.web.get_folder_by_server_relative_url(folder_path)
#         files = folder.expand(["Files"]).get().execute_query()
        
#         documents = []
#         for file in files.files:
#             with tempfile.NamedTemporaryFile(delete=False, dir=self.temp_dir) as tmp:
#                 file.download(tmp.name).execute_query()
                
#                 if file.name.endswith(('.mp4', '.mov', '.avi')):
#                     documents.extend(self._process_video(tmp.name, file))
#                 elif file.name.endswith(('.pptx', '.ppt')):
#                     documents.extend(self._process_slides(tmp.name, file))
#                 else:
#                     self.stats['other'] += 1
        
#         return documents

#     def _process_video(self, file_path, file_item):
#         # Transcribe with Whisper
#         transcript = self.llm_service.transcribe_audio(file_path)
        
#         # Split into time-coded chunks
#         return [Document(
#             text=segment['text'],
#             metadata={
#                 'type': 'video',
#                 'title': file_item.name,
#                 'timecodes': segment['start'],
#                 'sharepoint_url': file_item.serverRelativeUrl
#             }
#         ) for segment in transcript.segments]

#     def _process_slides(self, file_path, file_item):
#         presentation = pptx.Presentation(file_path)
#         return [Document(
#             text=f"Slide {i+1}: {slide.text}",
#             metadata={
#                 'type': 'slide',
#                 'title': file_item.name,
#                 'slide_number': i+1,
#                 'sharepoint_url': file_item.serverRelativeUrl
#             }
#         ) for i, slide in enumerate(presentation.slides)]

#     def get_processing_stats(self):
#         return self.stats